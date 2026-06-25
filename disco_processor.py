# -*- coding: ascii -*-
"""
SELLL.io -- Disco Call Processor
=================================
Feed it a Fathom transcript or summary and it outputs:
  - Prospect intelligence JSON
  - ICP score (0-100) with reasoning
  - Personalised "What We Heard" slide content
  - Follow-up email (ready to send)
  - CRM entry (paste-ready)
  - Personalized sales deck (.pptx)

USAGE (manual):
  python disco_processor.py --file fathom_notes.txt --company "Acme Corp"

USAGE (webhook server -- receives Fathom via n8n/Zapier):
  python disco_processor.py --serve
  Then POST JSON to http://localhost:5050/disco
  Body: { "transcript": "...", "company": "...", "date": "..." }

SETUP:
  Set ANTHROPIC_API_KEY in environment or create a .env file.
"""

import os
import sys
import json
import re
import argparse
import datetime
from pathlib import Path

# ------------------------------------------------------------------ paths --
BASE = Path(r"C:\Users\COLLINS OGIKI\Desktop\Extruct-GTM-SKILLS\SELLL-Engine")
OUT  = BASE / "Disco-Calls"
LOGO = Path(r"C:\Users\COLLIN~1\AppData\Local\Temp\claude\c--Users-COLLINS-OGIKI-Desktop-Extruct-GTM-SKILLS\d65f70c4-aefc-4b79-923f-049ab953d08e\scratchpad\selll_logo_clean.png")
OUT.mkdir(parents=True, exist_ok=True)

# ----------------------------------------------------------------- config --
MODEL   = "claude-haiku-4-5-20251001"   # fast + cheap for processing
MAX_TOK = 4096


# ===========================================================================
# 1. CLAUDE ANALYSIS ENGINE
# ===========================================================================

SYSTEM_PROMPT = """You are the SELLL.io Revenue Intelligence Engine.
SELLL builds signal-triggered outbound systems for B2B companies.
Your job is to analyse a discovery call transcript and extract
structured intelligence that powers personalised sales materials.

You must respond with ONLY valid JSON matching the schema exactly.
No prose, no markdown fences, just the raw JSON object."""

EXTRACTION_PROMPT = """Analyse this discovery call transcript and extract the following.
Return ONLY valid JSON with this exact structure (fill every field):

{
  "prospect": {
    "first_name": "",
    "last_name": "",
    "full_name": "",
    "title": "",
    "company": "",
    "email": "",
    "linkedin": "",
    "phone": ""
  },
  "company": {
    "name": "",
    "website": "",
    "headcount": "",
    "arr_estimate": "",
    "industry": "",
    "geography": "",
    "funding_stage": "",
    "founded_approx": "",
    "tech_stack": []
  },
  "pain_points": [
    {"label": "", "exact_quote": "", "severity": "Critical|High|Medium", "my_words": ""}
  ],
  "current_state": {
    "pipeline_description": "",
    "outbound_motion": "",
    "crm_tool": "",
    "email_tool": "",
    "has_sdr": false,
    "founder_selling": false,
    "referral_dependent": false,
    "tried_cold_outreach": false,
    "previous_agency": ""
  },
  "goals": {
    "primary_goal": "",
    "revenue_target": "",
    "timeline": "",
    "board_pressure": false,
    "funding_related_urgency": false
  },
  "signals": {
    "urgency_level": "Immediate|High|Medium|Low",
    "budget_signal": "Strong|Moderate|Weak|Unknown",
    "decision_authority": "Full|Partial|Influencer",
    "competition_mentioned": [],
    "objections_raised": [],
    "positive_buying_signals": []
  },
  "icp_score": {
    "total": 0,
    "firmographic": {"score": 0, "max": 25, "rationale": ""},
    "technographic": {"score": 0, "max": 15, "rationale": ""},
    "pain_alignment": {"score": 0, "max": 20, "rationale": ""},
    "budget_capacity": {"score": 0, "max": 20, "rationale": ""},
    "contact_access": {"score": 0, "max": 10, "rationale": ""},
    "timing_signals": {"score": 0, "max": 10, "rationale": ""},
    "grade": "A+|A|B|C|D",
    "summary": ""
  },
  "recommended_plan": {
    "tier": "Foundation|Growth|Enterprise",
    "monthly_price": "",
    "rationale": "",
    "first_three_signals": [],
    "suggested_sequences": []
  },
  "red_flags": [],
  "next_step": {
    "action": "",
    "deadline": "",
    "what_to_send": ""
  },
  "what_we_heard": {
    "pain_1": {"label": "", "detail": ""},
    "pain_2": {"label": "", "detail": ""},
    "pain_3": {"label": "", "detail": ""},
    "context": "",
    "urgency_statement": ""
  },
  "call_date": "",
  "call_summary": ""
}

TRANSCRIPT:
{transcript}
"""


def run_analysis(transcript: str, api_key: str) -> dict:
    import anthropic
    client = anthropic.Anthropic(api_key=api_key)
    prompt = EXTRACTION_PROMPT.replace("{transcript}", transcript[:12000])
    msg = client.messages.create(
        model=MODEL,
        max_tokens=MAX_TOK,
        system=SYSTEM_PROMPT,
        messages=[{"role": "user", "content": prompt}]
    )
    raw = msg.content[0].text.strip()
    # Strip any accidental markdown fences
    raw = re.sub(r"^```json\s*", "", raw)
    raw = re.sub(r"^```\s*",     "", raw)
    raw = re.sub(r"\s*```$",     "", raw)
    return json.loads(raw)


# ===========================================================================
# 2. OUTPUT GENERATORS
# ===========================================================================

def slug(text: str) -> str:
    return re.sub(r"[^a-z0-9]+", "-", text.lower()).strip("-")


def make_output_dir(intel: dict) -> Path:
    company  = intel.get("company", {}).get("name", "Unknown")
    date_str = intel.get("call_date", "") or datetime.date.today().isoformat()
    folder   = OUT / f"{slug(company)}-{date_str}"
    folder.mkdir(parents=True, exist_ok=True)
    return folder


def write_intel_json(folder: Path, intel: dict):
    path = folder / "prospect-intel.json"
    path.write_text(json.dumps(intel, indent=2), encoding="utf-8")
    return path


def write_icp_score(folder: Path, intel: dict) -> Path:
    s    = intel.get("icp_score", {})
    comp = intel.get("company", {})
    pros = intel.get("prospect", {})
    lines = [
        "=" * 60,
        "SELLL.io -- ICP SCORE CARD",
        "=" * 60,
        f"Prospect : {pros.get('full_name','')} -- {pros.get('title','')}",
        f"Company  : {comp.get('name','')}",
        f"Date     : {intel.get('call_date','')}",
        "",
        f"TOTAL SCORE : {s.get('total',0)} / 100   Grade: {s.get('grade','')}",
        "",
        f"  Firmographic   {s.get('firmographic',{}).get('score',0):>3} / 25  -- "
        f"{s.get('firmographic',{}).get('rationale','')}",
        f"  Technographic  {s.get('technographic',{}).get('score',0):>3} / 15  -- "
        f"{s.get('technographic',{}).get('rationale','')}",
        f"  Pain Alignment {s.get('pain_alignment',{}).get('score',0):>3} / 20  -- "
        f"{s.get('pain_alignment',{}).get('rationale','')}",
        f"  Budget         {s.get('budget_capacity',{}).get('score',0):>3} / 20  -- "
        f"{s.get('budget_capacity',{}).get('rationale','')}",
        f"  Contact Access {s.get('contact_access',{}).get('score',0):>3} / 10  -- "
        f"{s.get('contact_access',{}).get('rationale','')}",
        f"  Timing Signals {s.get('timing_signals',{}).get('score',0):>3} / 10  -- "
        f"{s.get('timing_signals',{}).get('rationale','')}",
        "",
        "SUMMARY:",
        s.get("summary", ""),
        "",
        "RED FLAGS:",
    ]
    for rf in intel.get("red_flags", []) or ["None identified."]:
        lines.append(f"  - {rf}")
    lines += [
        "",
        "RECOMMENDED PLAN:",
        f"  {intel.get('recommended_plan',{}).get('tier','')} at "
        f"{intel.get('recommended_plan',{}).get('monthly_price','')}",
        f"  {intel.get('recommended_plan',{}).get('rationale','')}",
    ]
    path = folder / "prospect-score.txt"
    path.write_text("\n".join(lines), encoding="utf-8")
    return path


def write_follow_up_email(folder: Path, intel: dict) -> Path:
    pros  = intel.get("prospect", {})
    comp  = intel.get("company",  {})
    wh    = intel.get("what_we_heard", {})
    ns    = intel.get("next_step", {})
    plan  = intel.get("recommended_plan", {})
    date  = intel.get("call_date", datetime.date.today().isoformat())
    pains = intel.get("pain_points", [])
    p1    = pains[0].get("exact_quote", pains[0].get("label","")) if len(pains) > 0 else ""
    p2    = pains[1].get("exact_quote", pains[1].get("label","")) if len(pains) > 1 else ""

    email = f"""Subject: recap from our call + what we would build for {comp.get('name','')}

Hi {pros.get('first_name','[FIRST NAME]')},

Thanks for the time today. A few things I took away from our conversation:

- "{p1}"
- "{p2}"
- {wh.get('context','')}
- {wh.get('urgency_statement','')}

Based on that, I am putting together a custom SELLL Build Plan for {comp.get('name','[COMPANY]')}.
It will show you:
  1. The signal infrastructure we would deploy for your specific ICP
  2. The 3 sequences we would build and exactly why
  3. Projected pipeline targets for months 1-3

I will have it to you by {ns.get('deadline','[DATE]')}.

One thing that will help me make it specific to {comp.get('name','[COMPANY]')}: {ns.get('what_to_send','')}

Recommended starting point for you: the {plan.get('tier','')} plan at {plan.get('monthly_price','')}/month.
{plan.get('rationale','')}

Talk soon,
Aaron
SELLL.io | selll.io/audit | hello@selll.io
"""
    path = folder / "follow-up-email.txt"
    path.write_text(email.strip(), encoding="utf-8")
    return path


def write_crm_entry(folder: Path, intel: dict) -> Path:
    pros  = intel.get("prospect", {})
    comp  = intel.get("company",  {})
    cs    = intel.get("current_state", {})
    sig   = intel.get("signals", {})
    score = intel.get("icp_score", {})
    plan  = intel.get("recommended_plan", {})
    ns    = intel.get("next_step", {})
    lines = [
        "=" * 60,
        "SELLL CRM ENTRY -- PASTE INTO HUBSPOT/PIPEDRIVE",
        "=" * 60,
        "",
        f"DEAL NAME     : {comp.get('name','')} -- Discovery",
        f"STAGE         : S4 -- Discovery Held",
        f"CLOSE DATE    : {ns.get('deadline','')}",
        f"DEAL VALUE    : {plan.get('monthly_price','')} x 3 months minimum",
        "",
        "CONTACT",
        f"  Name        : {pros.get('full_name','')}",
        f"  Title       : {pros.get('title','')}",
        f"  Email       : {pros.get('email','')}",
        f"  LinkedIn    : {pros.get('linkedin','')}",
        f"  Phone       : {pros.get('phone','')}",
        "",
        "COMPANY",
        f"  Name        : {comp.get('name','')}",
        f"  Website     : {comp.get('website','')}",
        f"  Headcount   : {comp.get('headcount','')}",
        f"  ARR (est.)  : {comp.get('arr_estimate','')}",
        f"  Industry    : {comp.get('industry','')}",
        f"  Funding     : {comp.get('funding_stage','')}",
        f"  Tech Stack  : {', '.join(comp.get('tech_stack',[]))}",
        "",
        "QUALIFICATION",
        f"  ICP Score   : {score.get('total',0)} / 100  Grade: {score.get('grade','')}",
        f"  Urgency     : {sig.get('urgency_level','')}",
        f"  Budget sig. : {sig.get('budget_signal','')}",
        f"  Authority   : {sig.get('decision_authority','')}",
        f"  Founder selling: {'Yes' if cs.get('founder_selling') else 'No'}",
        f"  Has SDR     : {'Yes' if cs.get('has_sdr') else 'No'}",
        f"  Tried outbound: {'Yes' if cs.get('tried_cold_outreach') else 'No'}",
        "",
        "KEY PAIN POINTS",
    ]
    for pp in intel.get("pain_points", []):
        lines.append(f"  [{pp.get('severity','')}] {pp.get('label','')} -- \"{pp.get('exact_quote','')}\"")
    lines += [
        "",
        "COMPETITION MENTIONED",
    ]
    comps = sig.get("competition_mentioned", []) or ["None mentioned"]
    for c in comps: lines.append(f"  - {c}")
    lines += [
        "",
        "OBJECTIONS RAISED",
    ]
    objs = sig.get("objections_raised", []) or ["None raised"]
    for o in objs: lines.append(f"  - {o}")
    lines += [
        "",
        "SIGNAL TYPE TAG : Discovery -- Direct outreach",
        f"RECOMMENDED PLAN: {plan.get('tier','')} at {plan.get('monthly_price','')}",
        "",
        "NEXT ACTION",
        f"  {ns.get('action','')}",
        f"  Deadline: {ns.get('deadline','')}",
        f"  Send: {ns.get('what_to_send','')}",
        "",
        f"NOTES: {intel.get('call_summary','')}",
    ]
    path = folder / "crm-entry.txt"
    path.write_text("\n".join(lines), encoding="utf-8")
    return path


def write_what_we_heard(folder: Path, intel: dict) -> Path:
    wh    = intel.get("what_we_heard", {})
    comp  = intel.get("company", {})
    pros  = intel.get("prospect", {})
    date  = intel.get("call_date", "")
    lines = [
        "=" * 60,
        "SLIDE 2 -- WHAT WE HEARD",
        f"Prepared for: {comp.get('name','')}  |  Call: {date}",
        "=" * 60,
        "",
        "REPLACE THE [BRACKETS] IN YOUR SALES DECK SLIDE 2 WITH:",
        "",
        f"PAIN 1 -- {wh.get('pain_1',{}).get('label','')}",
        f"  Detail: {wh.get('pain_1',{}).get('detail','')}",
        "",
        f"PAIN 2 -- {wh.get('pain_2',{}).get('label','')}",
        f"  Detail: {wh.get('pain_2',{}).get('detail','')}",
        "",
        f"PAIN 3 -- {wh.get('pain_3',{}).get('label','')}",
        f"  Detail: {wh.get('pain_3',{}).get('detail','')}",
        "",
        f"CONTEXT: {wh.get('context','')}",
        "",
        f"URGENCY: {wh.get('urgency_statement','')}",
        "",
        "=" * 60,
        "DECK PERSONALISATION CHECKLIST",
        "  Slide 1:  Add company name in the top-right tag",
        f"  Slide 2:  Replace all [BRACKETS] using the above",
        f"  Slide 7:  Fill in [COMPANY] signals with:",
    ]
    for sig in intel.get("recommended_plan", {}).get("first_three_signals", []):
        lines.append(f"    - {sig}")
    lines += [
        f"  Slide 10: Replace [X,XXX] ACV with: "
        f"{intel.get('company',{}).get('arr_estimate','')} estimated ARR",
        f"  Slide 11: Highlight the {intel.get('recommended_plan',{}).get('tier','')} plan",
        "=" * 60,
    ]
    path = folder / "what-we-heard.txt"
    path.write_text("\n".join(lines), encoding="utf-8")
    return path


def build_personalized_deck(folder: Path, intel: dict):
    """Builds a personalized SELLL-Sales-Deck with the prospect's info filled in."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        DARK = RGBColor(0x1A, 0x1A, 0x2E)
        GOLD = RGBColor(0xE8, 0xB8, 0x30)
        WHITE = RGBColor(0xFF, 0xFF, 0xFF)
        GREY  = RGBColor(0x88, 0x88, 0x99)
        LGREY = RGBColor(0xCC, 0xCC, 0xD8)
        MID   = RGBColor(0x26, 0x26, 0x40)
        NAVY2 = RGBColor(0x0F, 0x0F, 0x1C)

        pros  = intel.get("prospect", {})
        comp  = intel.get("company",  {})
        wh    = intel.get("what_we_heard", {})
        plan  = intel.get("recommended_plan", {})
        score = intel.get("icp_score", {})
        date  = intel.get("call_date", datetime.date.today().isoformat())
        pains = intel.get("pain_points", [])

        prs = Presentation()
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)

        def blank(): return prs.slides.add_slide(prs.slide_layouts[6])
        def bgc(sl, rgb):
            f = sl.background.fill; f.solid(); f.fore_color.rgb = rgb
        def R(sl, l, t, w, h, rgb):
            s = sl.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
            s.fill.solid(); s.fill.fore_color.rgb = rgb
            s.line.fill.background(); return s
        def T(sl, text, l, t, w, h, sz, bold=False, col=WHITE, al=PP_ALIGN.LEFT, it=False):
            box = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
            tf = box.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.alignment = al
            run = p.add_run(); run.text = text
            run.font.size = Pt(sz); run.font.bold = bold
            run.font.color.rgb = col; run.font.italic = it
        def gold_rule(sl, x, y, length):
            R(sl, x, y, length, 0.03, GOLD)
        def snum(sl, n):
            T(sl, str(n), 12.82, 7.08, 0.46, 0.36, 9, col=GREY, al=PP_ALIGN.RIGHT)

        company_name = comp.get("name", "[COMPANY]")
        fname        = pros.get("first_name", "[FIRST NAME]")
        score_total  = score.get("total", 0)
        grade        = score.get("grade", "")

        # ---- SLIDE 1: Cover ----
        sl = blank(); bgc(sl, NAVY2)
        R(sl, 0, 0, 13.33, 7.5, DARK)
        R(sl, 0, 0, 0.06, 7.5, GOLD); R(sl, 13.27, 0, 0.06, 7.5, GOLD)
        R(sl, 0, 0, 13.33, 0.06, GOLD); R(sl, 0, 7.44, 13.33, 0.06, GOLD)
        R(sl, 9.2, 0.06, 4.07, 0.62, GOLD)
        T(sl, "Prepared for:  " + company_name, 9.32, 0.16, 3.9, 0.38, 11, bold=True, col=DARK)
        if LOGO.exists():
            sl.shapes.add_picture(str(LOGO), Inches(0.38), Inches(0.28), width=Inches(3.2))
        T(sl, "Your Pipeline is Predictable.", 0.55, 1.6, 11.5, 1.0, 44, bold=True, col=WHITE)
        T(sl, "We build the system that makes it that way.", 0.55, 2.58, 11.5, 0.7, 28, col=GOLD)
        gold_rule(sl, 0.55, 3.42, 6.5)
        for i, (icon, lbl) in enumerate([
            ("30 days", "to first qualified meeting"),
            ("90 days", "to a fully documented system"),
            ("Day 1",   "full IP ownership - yours forever"),
        ]):
            x = 0.55 + i * 4.2
            R(sl, x, 3.6, 3.9, 0.98, MID); R(sl, x, 3.6, 3.9, 0.05, GOLD)
            T(sl, icon, x+0.2, 3.72, 3.5, 0.42, 24, bold=True, col=GOLD)
            T(sl, lbl,  x+0.2, 4.14, 3.5, 0.34, 9,  col=LGREY)
        T(sl, "SELLL.io  |  Signal-Intelligent Revenue Engine  |  selll.io/audit",
          0.55, 7.0, 12.2, 0.38, 9, col=GREY, al=PP_ALIGN.CENTER)
        snum(sl, 1)

        # ---- SLIDE 2: What We Heard (FILLED IN) ----
        sl = blank(); bgc(sl, RGBColor(0xFF, 0xFC, 0xF0))
        R(sl, 0, 0, 13.33, 0.055, GOLD)
        T(sl, "BEFORE WE START", 0.55, 0.22, 6.0, 0.38, 9, bold=True,
          col=RGBColor(0xB8, 0x86, 0x0B))
        T(sl, "Here Is What We Heard\nOn Your Call.", 0.55, 0.72, 7.5, 1.3, 34, bold=True, col=DARK)
        T(sl, "This deck was built for " + company_name + " based on your conversation on " + date + ".",
          0.55, 2.1, 8.2, 0.5, 11, it=True, col=GREY)
        gold_rule(sl, 0.55, 2.72, 5.5)
        heard_data = [
            (wh.get("pain_1", {}).get("label", ""),  wh.get("pain_1", {}).get("detail", "")),
            (wh.get("pain_2", {}).get("label", ""),  wh.get("pain_2", {}).get("detail", "")),
            (wh.get("context", "Context"),             wh.get("context", "")),
            (wh.get("urgency_statement", "Urgency"),  wh.get("urgency_statement", "")),
        ]
        for i, (label, detail) in enumerate(heard_data):
            col_ = i % 2; row_ = i // 2
            x = 0.48 + col_ * 6.42; y = 2.92 + row_ * 2.05
            R(sl, x, y, 6.18, 1.82, DARK); R(sl, x, y, 6.18, 0.05, GOLD)
            T(sl, label,  x+0.2, y+0.12, 5.8, 0.35, 9, bold=True, col=GOLD)
            T(sl, '"' + detail[:180] + '"', x+0.2, y+0.52, 5.7, 1.15, 10, it=True, col=LGREY)
        snum(sl, 2)

        # ---- SLIDE 3: ICP Score Summary ----
        sl = blank(); bgc(sl, DARK)
        R(sl, 0, 0, 13.33, 0.055, GOLD); R(sl, 0, 0, 0.07, 7.5, GOLD)
        T(sl, "PROSPECT INTELLIGENCE", 0.55, 0.22, 6.0, 0.38, 9, bold=True, col=GOLD)
        T(sl, company_name + " -- ICP Score: " + str(score_total) + "/100",
          0.55, 0.72, 10.0, 0.82, 34, bold=True, col=WHITE)
        T(sl, "Grade: " + grade + "  |  " + score.get("summary", ""),
          0.55, 1.65, 10.5, 0.55, 13, col=LGREY, it=True)
        gold_rule(sl, 0.55, 2.32, 12.3)
        cats = [
            ("Firmographic",   score.get("firmographic",  {}).get("score",0), 25),
            ("Technographic",  score.get("technographic", {}).get("score",0), 15),
            ("Pain Alignment", score.get("pain_alignment",{}).get("score",0), 20),
            ("Budget",         score.get("budget_capacity",{}).get("score",0),20),
            ("Contact Access", score.get("contact_access",{}).get("score",0), 10),
            ("Timing Signals", score.get("timing_signals",{}).get("score",0), 10),
        ]
        for i, (cat, sc, mx) in enumerate(cats):
            col_ = i % 2; row_ = i // 2
            x = 0.48 + col_ * 6.42; y = 2.52 + row_ * 1.55
            R(sl, x, y, 6.18, 1.28, MID); R(sl, x, y, 6.18, 0.05, GOLD)
            T(sl, cat, x+0.2, y+0.12, 3.5, 0.35, 10, bold=True, col=WHITE)
            T(sl, str(sc) + " / " + str(mx), x+4.5, y+0.12, 1.55, 0.35, 14, bold=True,
              col=GOLD, al=PP_ALIGN.RIGHT)
            bar_w = max(0.05, (sc / mx) * 5.5)
            R(sl, x+0.2, y+0.62, 5.5, 0.22, RGBColor(0x33, 0x33, 0x50))
            R(sl, x+0.2, y+0.62, bar_w, 0.22, GOLD)
            rationale = score.get(cat.lower().replace(" ", "_"), {}).get("rationale", "")
            T(sl, rationale[:90], x+0.2, y+0.92, 5.7, 0.28, 7, col=LGREY, it=True)
        T(sl, "Recommended plan: " + plan.get("tier","") + " at " + plan.get("monthly_price",""),
          0.55, 7.08, 8.0, 0.3, 9, bold=True, col=GOLD)
        snum(sl, 3)

        # ---- SLIDE 4: Next Steps ----
        sl = blank(); bgc(sl, NAVY2)
        R(sl, 0, 0, 13.33, 0.06, GOLD)
        R(sl, 0, 0, 0.06, 7.5, GOLD); R(sl, 13.27, 0, 0.06, 7.5, GOLD)
        R(sl, 0, 7.44, 13.33, 0.06, GOLD)
        T(sl, "WHAT HAPPENS NEXT", 1.5, 0.3, 10.3, 0.42, 11, bold=True, col=GOLD, al=PP_ALIGN.CENTER)
        T(sl, "Three Steps to Your First\nQualified Meeting.",
          1.2, 0.88, 10.9, 1.32, 36, bold=True, col=WHITE, al=PP_ALIGN.CENTER)
        gold_rule(sl, 4.5, 2.75, 4.3)
        ns = intel.get("next_step", {})
        steps = [
            ("01", "Say yes to the Build Plan",
             "Reply to my follow-up email. We send the agreement, you sign, "
             "we start building Day 1. No delay. No onboarding forms. We work."),
            ("02", "30-minute kickoff call",
             "Within 48 hours of signing we hold a kickoff. We confirm your ICP, "
             "signal priorities, and get CRM + domain access. "
             "By end of the call, your system build has started."),
            ("03", "First qualified meeting in 30 days",
             "Sequences fire by Day 28-30. Replies arrive. "
             "Meetings land in your calendar. You show up and close."),
        ]
        for i, (num, step, body) in enumerate(steps):
            y = 3.0 + i * 1.3
            R(sl, 1.82, y, 0.6, 0.6, GOLD)
            T(sl, num, 1.84, y+0.08, 0.56, 0.44, 18, bold=True, col=DARK, al=PP_ALIGN.CENTER)
            T(sl, step, 2.58, y+0.04, 8.2, 0.38, 13, bold=True, col=WHITE)
            T(sl, body, 2.58, y+0.46, 8.2, 0.7,  9,  col=LGREY)
        R(sl, 2.8, 7.0, 7.7, 0.46, GOLD)
        T(sl, "Ready to start?  Reply to this email or visit  selll.io/audit",
          2.92, 7.06, 7.46, 0.32, 12, bold=True, col=DARK, al=PP_ALIGN.CENTER)
        snum(sl, 4)

        out_path = folder / f"SELLL-Sales-Deck-{slug(company_name)}.pptx"
        prs.save(str(out_path))
        return out_path

    except Exception as e:
        return None


# ===========================================================================
# 3. MAIN PROCESSOR
# ===========================================================================

def process(transcript: str, api_key: str) -> dict:
    print("\n[SELLL] Analysing transcript with Claude...")
    intel  = run_analysis(transcript, api_key)
    folder = make_output_dir(intel)
    comp   = intel.get("company", {}).get("name", "Unknown")
    grade  = intel.get("icp_score", {}).get("grade", "?")
    score  = intel.get("icp_score", {}).get("total", 0)

    print(f"[SELLL] Company : {comp}")
    print(f"[SELLL] Grade   : {grade}  ({score}/100)")

    write_intel_json(folder, intel)
    print(f"[SELLL] Intel JSON     -> {folder / 'prospect-intel.json'}")

    write_icp_score(folder, intel)
    print(f"[SELLL] Score card     -> {folder / 'prospect-score.txt'}")

    write_what_we_heard(folder, intel)
    print(f"[SELLL] What we heard  -> {folder / 'what-we-heard.txt'}")

    write_follow_up_email(folder, intel)
    print(f"[SELLL] Follow-up email-> {folder / 'follow-up-email.txt'}")

    write_crm_entry(folder, intel)
    print(f"[SELLL] CRM entry      -> {folder / 'crm-entry.txt'}")

    deck_path = build_personalized_deck(folder, intel)
    if deck_path:
        print(f"[SELLL] Sales deck     -> {deck_path}")

    print(f"\n[SELLL] All outputs saved to: {folder}")
    return intel


# ===========================================================================
# 4. WEBHOOK SERVER (receives Fathom via n8n / Zapier / direct)
# ===========================================================================

def run_server(api_key: str, port: int = 5050):
    from flask import Flask, request, jsonify
    app = Flask(__name__)

    @app.route("/disco", methods=["POST"])
    def disco():
        data       = request.get_json(force=True) or {}
        transcript = (data.get("transcript") or data.get("summary") or
                      data.get("notes") or data.get("text") or "")
        if not transcript:
            return jsonify({"error": "No transcript provided"}), 400
        try:
            intel = process(transcript, api_key)
            return jsonify({
                "status":  "ok",
                "company": intel.get("company", {}).get("name"),
                "grade":   intel.get("icp_score", {}).get("grade"),
                "score":   intel.get("icp_score", {}).get("total"),
                "folder":  str(OUT / (
                    slug(intel.get("company", {}).get("name","unknown")) +
                    "-" + intel.get("call_date", datetime.date.today().isoformat())
                )),
            })
        except Exception as e:
            return jsonify({"error": str(e)}), 500

    @app.route("/health", methods=["GET"])
    def health():
        return jsonify({"status": "SELLL disco processor online"})

    print(f"\n[SELLL] Webhook server running on http://localhost:{port}")
    print(f"[SELLL] POST your Fathom transcript to http://localhost:{port}/disco")
    print(f"[SELLL] Body: {{\"transcript\": \"...\"}}")
    print(f"[SELLL] Health check: http://localhost:{port}/health\n")
    app.run(host="0.0.0.0", port=port, debug=False)


# ===========================================================================
# 5. CLI ENTRY POINT
# ===========================================================================

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="SELLL Disco Call Processor")
    parser.add_argument("--file",    help="Path to Fathom transcript .txt file")
    parser.add_argument("--text",    help="Paste transcript directly as a string")
    parser.add_argument("--serve",   action="store_true",
                        help="Run webhook server (port 5050)")
    parser.add_argument("--port",    type=int, default=5050)
    args = parser.parse_args()

    # Load API key
    api_key = os.environ.get("ANTHROPIC_API_KEY", "")
    if not api_key:
        env_file = Path(__file__).parent / ".env"
        if env_file.exists():
            for line in env_file.read_text().splitlines():
                if line.startswith("ANTHROPIC_API_KEY="):
                    api_key = line.split("=", 1)[1].strip().strip('"')
    if not api_key:
        print("[ERROR] ANTHROPIC_API_KEY not found.")
        print("  Option 1: set environment variable ANTHROPIC_API_KEY=sk-ant-...")
        print("  Option 2: create a .env file next to this script with:")
        print("            ANTHROPIC_API_KEY=sk-ant-...")
        sys.exit(1)

    if args.serve:
        run_server(api_key, args.port)

    elif args.file:
        transcript = Path(args.file).read_text(encoding="utf-8", errors="ignore")
        process(transcript, api_key)

    elif args.text:
        process(args.text, api_key)

    else:
        print("SELLL Disco Processor -- Usage:")
        print("  Process a file:    python disco_processor.py --file fathom_notes.txt")
        print("  Paste transcript:  python disco_processor.py --text \"full transcript here\"")
        print("  Run webhook:       python disco_processor.py --serve")
        print("")
        print("Set ANTHROPIC_API_KEY in environment or .env file first.")
